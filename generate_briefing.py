import os
import json
import requests
from datetime import datetime
from pptx import Presentation
from io import BytesIO
import anthropic

# ── SharePoint 인증 ──────────────────────────────
CLIENT_ID = os.environ['SP_CLIENT_ID']
CLIENT_SECRET = os.environ['SP_CLIENT_SECRET']
TENANT_ID = os.environ['SP_TENANT_ID']
CLAUDE_API_KEY = os.environ['CLAUDE_API_KEY']

SITE_URL = "https://enjetco.sharepoint.com/sites/98"
LIST_NAME = "경영회의DB"

def get_sp_token():
    url = f"https://login.microsoftonline.com/{TENANT_ID}/oauth2/v2.0/token"
    data = {
        "client_id": CLIENT_ID,
        "client_secret": CLIENT_SECRET,
        "scope": "https://graph.microsoft.com/.default",
        "grant_type": "client_credentials"
    }
    r = requests.post(url, data=data)
    return r.json()["access_token"]

def get_sp_token_rest():
    url = f"https://accounts.accesscontrol.windows.net/{TENANT_ID}/tokens/OAuth/2"
    data = {
        "grant_type": "client_credentials",
        "client_id": f"{CLIENT_ID}@{TENANT_ID}",
        "client_secret": CLIENT_SECRET,
        "resource": f"00000003-0000-0ff1-ce00-000000000000/enjetco.sharepoint.com@{TENANT_ID}"
    }
    r = requests.post(url, data=data)
    return r.json().get("access_token")

def get_pending_items(token):
    """경영회의DB에서 AI요약이 없는 항목 가져오기"""
    headers = {
        "Authorization": f"Bearer {token}",
        "Accept": "application/json"
    }
    # Graph API로 SharePoint List 조회
    site_id_url = f"https://graph.microsoft.com/v1.0/sites/enjetco.sharepoint.com:/sites/98"
    site_r = requests.get(site_id_url, headers=headers)
    site_id = site_r.json().get("id", "")

    list_url = f"https://graph.microsoft.com/v1.0/sites/{site_id}/lists/{LIST_NAME}/items?expand=fields&$filter=fields/AI요약 eq null"
    r = requests.get(list_url, headers=headers)
    return r.json().get("value", [])

def download_ppt(token, file_path):
    """SharePoint에서 PPT 다운로드"""
    headers = {"Authorization": f"Bearer {token}"}
    # 파일 경로로 다운로드
    encoded_path = requests.utils.quote(file_path)
    url = f"https://enjetco.sharepoint.com/sites/98/_api/web/GetFileByServerRelativeUrl('{file_path}')/$value"
    
    # REST API 토큰으로 다운로드
    rest_token = get_sp_token_rest()
    headers_rest = {"Authorization": f"Bearer {rest_token}"}
    r = requests.get(url, headers=headers_rest)
    return r.content

def extract_ppt_text(ppt_bytes):
    """PPT에서 텍스트 추출"""
    prs = Presentation(BytesIO(ppt_bytes))
    text = []
    for i, slide in enumerate(prs.slides):
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                slide_texts.append(shape.text.strip())
        if slide_texts:
            text.append(f"[슬라이드 {i+1}]\n" + "\n".join(slide_texts))
    return "\n\n".join(text)

def summarize_with_claude(ppt_text, dept_name):
    """Claude API로 요약"""
    client = anthropic.Anthropic(api_key=CLAUDE_API_KEY)
    
    prompt = f"""당신은 엔젯㈜ 경영회의 AI 어시스턴트입니다.
아래는 {dept_name} 주간 경영회의 PPT 내용입니다.

{ppt_text}

다음 형식으로 분석해주세요:
1. 핵심 안건 3가지 (각 2줄 이내)
2. 주요 수치/일정 (날짜, 목표치 등)
3. 리스크 또는 결정 필요 사항
4. TO DO 항목 (담당자 포함)

간결하고 명확하게 작성해주세요."""

    message = client.messages.create(
        model="claude-sonnet-4-20250514",
        max_tokens=1000,
        messages=[{"role": "user", "content": prompt}]
    )
    return message.content[0].text

def generate_briefing_html(summaries, week):
    """브리핑 HTML 생성"""
    dept_cards = ""
    for s in summaries:
        dept_cards += f"""
        <div class="dept-card">
            <div class="dept-header">
                <div class="dept-name">{s['dept']}</div>
            </div>
            <div class="dept-body">
                <pre class="summary">{s['summary']}</pre>
            </div>
        </div>"""

    html = f"""<!DOCTYPE html>
<html lang="ko">
<head>
<meta charset="UTF-8">
<title>{week} 경영회의 브리핑</title>
<style>
body {{ font-family: 'Noto Sans KR', sans-serif; background:#0A0F1A; color:#E2EAF8; margin:0; padding:24px; }}
h1 {{ text-align:center; color:#3D8BFF; margin-bottom:8px; }}
.subtitle {{ text-align:center; color:#5A7A99; margin-bottom:32px; font-size:13px; }}
.grid {{ display:grid; grid-template-columns:repeat(2,1fr); gap:16px; max-width:1200px; margin:0 auto; }}
.dept-card {{ background:#111827; border:1px solid #1E3050; border-radius:12px; overflow:hidden; }}
.dept-header {{ padding:14px 18px; background:#162035; border-bottom:1px solid #1E3050; }}
.dept-name {{ font-size:14px; font-weight:700; color:#3D8BFF; }}
.dept-body {{ padding:16px 18px; }}
.summary {{ font-size:12px; line-height:1.8; white-space:pre-wrap; color:#E2EAF8; margin:0; }}
.generated {{ text-align:center; color:#5A7A99; font-size:11px; margin-top:24px; }}
</style>
</head>
<body>
<h1>🤖 {week} 주간 경영회의 AI 브리핑</h1>
<div class="subtitle">Claude AI 자동 분석 · {datetime.now().strftime('%Y.%m.%d %H:%M')} 생성</div>
<div class="grid">
{dept_cards}
</div>
<div class="generated">이 브리핑은 SharePoint PPT 업로드 후 GitHub Actions + Claude AI로 자동 생성됩니다</div>
</body>
</html>"""
    return html

def update_sp_item(token, site_id, item_id, summary):
    """SP List 항목 AI요약 업데이트"""
    headers = {
        "Authorization": f"Bearer {token}",
        "Content-Type": "application/json"
    }
    url = f"https://graph.microsoft.com/v1.0/sites/{site_id}/lists/{LIST_NAME}/items/{item_id}/fields"
    data = {"AI요약": summary[:500]}
    requests.patch(url, headers=headers, json=data)

def main():
    print("🚀 브리핑 생성 시작...")
    
    token = get_sp_token()
    
    # SP List에서 새 항목 가져오기
    items = get_pending_items(token)
    print(f"📋 처리할 항목: {len(items)}개")
    
    if not items:
        print("새 PPT 없음. 종료.")
        return

    summaries = []
    week = datetime.now().strftime("W%V")

    for item in items:
        fields = item.get("fields", {})
        file_path = fields.get("PPT링크", "")
        title = fields.get("Title", "")
        dept = fields.get("발표부서", title)
        
        print(f"📊 처리중: {dept}")
        
        try:
            # PPT 다운로드
            ppt_bytes = download_ppt(token, file_path)
            # 텍스트 추출
            ppt_text = extract_ppt_text(ppt_bytes)
            # Claude 요약
            summary = summarize_with_claude(ppt_text, dept)
            summaries.append({"dept": dept, "summary": summary})
            print(f"✅ {dept} 요약 완료")
        except Exception as e:
            print(f"❌ {dept} 오류: {e}")
            summaries.append({"dept": dept, "summary": f"처리 오류: {str(e)}"})

    # HTML 생성
    html = generate_briefing_html(summaries, week)
    
    # 파일 저장
    os.makedirs("briefing", exist_ok=True)
    filename = f"briefing/{week}.html"
    with open(filename, "w", encoding="utf-8") as f:
        f.write(html)
    
    # index 페이지도 업데이트
    with open("briefing/latest.html", "w", encoding="utf-8") as f:
        f.write(html)
    
    print(f"✅ 브리핑 생성 완료: {filename}")

if __name__ == "__main__":
    main()
